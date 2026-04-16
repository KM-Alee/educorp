from __future__ import annotations

import io

import pdfplumber
import srt
import webvtt
from docx import Document
from pptx import Presentation

from app.services.canonical_builder import CanonicalPageRecord, build_canonical_page
from app.services.ocr_service import OCRService
from app.services.pdf_extraction import PDFExtractionPipeline
from app.services.vision_service import VisionService


class ExtractionService:
    """
    Orchestrates per-asset-type extraction into canonical page records.

    PDF assets use the full staged pipeline (native text → heuristics → OCR → vision).
    PPTX assets use slide-level extraction (text + speaker notes).
    All other formats (txt, md, docx, vtt, srt) use local-only extraction.
    """

    def __init__(self) -> None:
        self._ocr = OCRService()
        self._vision = VisionService()
        self._pdf_pipeline = PDFExtractionPipeline(
            ocr_service=self._ocr,
            vision_service=self._vision,
        )

    async def extract_canonical_pages(
        self,
        *,
        asset_type: str,
        data: bytes,
        asset_id: str,
        version_id: str,
        module_id: str,
        module_title: str,
        asset_title: str,
    ) -> tuple[list[CanonicalPageRecord], dict[str, int]]:
        """
        Extract a list of canonical page records and extraction stats.

        Stats keys: ``total_pages``, ``ocr_pages``, ``nanogpt_pages``,
        ``low_confidence_pages``.
        """
        if asset_type == "pdf":
            return await self._extract_pdf(
                data=data,
                asset_id=asset_id,
                version_id=version_id,
                module_id=module_id,
                module_title=module_title,
                asset_title=asset_title,
            )

        if asset_type == "pptx":
            return self._extract_pptx(
                data=data,
                asset_id=asset_id,
                version_id=version_id,
                module_id=module_id,
                module_title=module_title,
                asset_title=asset_title,
            )

        # All other formats: single-page canonical record
        text = _extract_flat(asset_type, data)
        record = build_canonical_page(
            version_id=version_id,
            asset_id=asset_id,
            source_type=asset_type,
            page_or_slide_number=1,
            module_id=module_id,
            module_title=module_title,
            asset_title=asset_title,
            native_text=text,
            text_confidence=1.0,
        )
        stats: dict[str, int] = {
            "total_pages": 1,
            "ocr_pages": 0,
            "nanogpt_pages": 0,
            "low_confidence_pages": 0,
        }
        return [record], stats

    # ── Private per-format helpers ────────────────────────────────────────────

    async def _extract_pdf(
        self,
        *,
        data: bytes,
        asset_id: str,
        version_id: str,
        module_id: str,
        module_title: str,
        asset_title: str,
    ) -> tuple[list[CanonicalPageRecord], dict[str, int]]:
        pages, stats = await self._pdf_pipeline.extract(asset_id=asset_id, data=data)
        records: list[CanonicalPageRecord] = [
            build_canonical_page(
                version_id=version_id,
                asset_id=asset_id,
                source_type="pdf",
                page_or_slide_number=page.page_number,
                module_id=module_id,
                module_title=module_title,
                asset_title=asset_title,
                native_text=page.native_text,
                ocr_text=page.ocr_text,
                visual_summary=page.visual_summary,
                has_visual_summary=page.has_visual_summary,
                text_confidence=page.text_confidence,
                visual_confidence=page.visual_confidence,
                flags=page.flags,
            )
            for page in pages
        ]
        return records, stats

    def _extract_pptx(
        self,
        *,
        data: bytes,
        asset_id: str,
        version_id: str,
        module_id: str,
        module_title: str,
        asset_title: str,
    ) -> tuple[list[CanonicalPageRecord], dict[str, int]]:
        presentation = Presentation(io.BytesIO(data))
        records: list[CanonicalPageRecord] = []
        for slide_num, slide in enumerate(presentation.slides, start=1):
            text_parts: list[str] = []
            notes_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text.strip())
            if slide.has_notes_slide:
                notes_tf = slide.notes_slide.notes_text_frame
                if notes_tf and notes_tf.text.strip():
                    notes_text = notes_tf.text.strip()
            native = "\n".join(text_parts)
            if notes_text:
                native = f"{native}\n\nNotes: {notes_text}" if native else notes_text
            records.append(
                build_canonical_page(
                    version_id=version_id,
                    asset_id=asset_id,
                    source_type="pptx",
                    page_or_slide_number=slide_num,
                    module_id=module_id,
                    module_title=module_title,
                    asset_title=asset_title,
                    native_text=native,
                    text_confidence=1.0 if native.strip() else 0.3,
                )
            )
        stats: dict[str, int] = {
            "total_pages": len(records),
            "ocr_pages": 0,
            "nanogpt_pages": 0,
            "low_confidence_pages": sum(
                1 for r in records if r.text_confidence < 0.65
            ),
        }
        return records, stats


# ── Flat / single-page extractors ────────────────────────────────────────────


def _extract_flat(asset_type: str, data: bytes) -> str:
    if asset_type in {"txt", "md"}:
        return data.decode("utf-8", errors="replace")
    if asset_type == "pdf":
        return _extract_pdf_flat(data)
    if asset_type == "docx":
        return _extract_docx(data)
    if asset_type == "vtt":
        return _extract_vtt(data)
    if asset_type == "srt":
        return _extract_srt(data)
    return data.decode("utf-8", errors="replace")


def _extract_pdf_flat(data: bytes) -> str:
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        pages = [page.extract_text() or "" for page in pdf.pages]
    return "\n".join(pages)


def _extract_docx(data: bytes) -> str:
    doc = Document(io.BytesIO(data))
    return "\n".join(para.text for para in doc.paragraphs)


def _extract_vtt(data: bytes) -> str:
    text = data.decode("utf-8", errors="replace")
    try:
        captions = webvtt.read_buffer(io.StringIO(text))
        return "\n".join(cap.text for cap in captions)
    except Exception:
        return text


def _extract_srt(data: bytes) -> str:
    text = data.decode("utf-8", errors="replace")
    try:
        subtitles = list(srt.parse(text))
        return "\n".join(sub.content for sub in subtitles)
    except Exception:
        return text
